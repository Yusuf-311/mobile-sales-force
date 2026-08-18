from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ──────────────────────────────────────────────────────────────────
TEAL      = RGBColor(0x00, 0x8B, 0x8B)
TEAL_DARK = RGBColor(0x00, 0x5F, 0x6B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG   = RGBColor(0xF4, 0xF7, 0xF9)
GRAY_TEXT = RGBColor(0x44, 0x55, 0x66)
ORANGE    = RGBColor(0xFF, 0x6B, 0x35)
GREEN     = RGBColor(0x2E, 0xCC, 0x71)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def title_slide(title, subtitle):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, TEAL_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, ORANGE)
    add_rect(slide, 0, 7.42, 13.33, 0.08, ORANGE)
    add_text(slide, title,    0.5, 2.5, 12, 1.5, size=40, bold=True,  color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, 0.5, 4.2, 12, 0.8, size=20, bold=False, color=RGBColor(0xCC,0xEE,0xFF), align=PP_ALIGN.CENTER)
    add_text(slide, "PT Mersifarma — Posisi: Software Engineer", 0.5, 5.2, 12, 0.5, size=14, color=RGBColor(0xAA,0xCC,0xDD), align=PP_ALIGN.CENTER)
    return slide

def section_slide(label):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, GRAY_BG)
    add_rect(slide, 0, 3.0, 13.33, 1.5, TEAL)
    add_text(slide, label, 0.5, 3.15, 12, 1.2, size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return slide

def content_slide(title, bullets, icon="▪"):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, GRAY_BG)
    add_rect(slide, 0, 0, 13.33, 1.2, TEAL)
    add_rect(slide, 0, 1.2, 0.06, 6.3, ORANGE)
    add_text(slide, title, 0.3, 0.2, 12, 0.85, size=26, bold=True, color=WHITE)
    y = 1.4
    for b in bullets:
        if b.startswith("##"):
            add_text(slide, b[2:].strip(), 0.4, y, 12.5, 0.4, size=14, bold=True, color=TEAL_DARK)
            y += 0.45
        else:
            add_text(slide, f"{icon}  {b}", 0.5, y, 12.3, 0.38, size=13, color=GRAY_TEXT)
            y += 0.42
    return slide

def two_col_slide(title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, GRAY_BG)
    add_rect(slide, 0, 0, 13.33, 1.2, TEAL)
    add_rect(slide, 0, 1.2, 0.06, 6.3, ORANGE)
    add_text(slide, title, 0.3, 0.2, 12, 0.85, size=26, bold=True, color=WHITE)
    # Left col
    add_rect(slide, 0.3, 1.35, 6.0, 5.8, WHITE)
    add_text(slide, left_title, 0.4, 1.45, 5.8, 0.5, size=15, bold=True, color=TEAL_DARK)
    y = 2.05
    for item in left_items:
        add_text(slide, f"▪  {item}", 0.5, y, 5.6, 0.38, size=12, color=GRAY_TEXT)
        y += 0.42
    # Right col
    add_rect(slide, 6.7, 1.35, 6.3, 5.8, WHITE)
    add_text(slide, right_title, 6.8, 1.45, 6.1, 0.5, size=15, bold=True, color=TEAL_DARK)
    y = 2.05
    for item in right_items:
        add_text(slide, f"▪  {item}", 6.9, y, 5.9, 0.38, size=12, color=GRAY_TEXT)
        y += 0.42
    return slide

def flow_slide(title, steps):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, GRAY_BG)
    add_rect(slide, 0, 0, 13.33, 1.2, TEAL)
    add_rect(slide, 0, 1.2, 0.06, 6.3, ORANGE)
    add_text(slide, title, 0.3, 0.2, 12, 0.85, size=26, bold=True, color=WHITE)
    n = len(steps)
    box_w = 11.5 / n - 0.15
    for i, (label, desc) in enumerate(steps):
        x = 0.7 + i * (box_w + 0.2)
        add_rect(slide, x, 2.0, box_w, 1.4, TEAL)
        add_text(slide, label, x+0.05, 2.1, box_w-0.1, 0.6, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, x, 3.5, box_w, 2.5, WHITE)
        add_text(slide, desc, x+0.08, 3.6, box_w-0.16, 2.3, size=11, color=GRAY_TEXT, wrap=True)
        if i < n-1:
            add_text(slide, "→", x+box_w+0.01, 2.45, 0.2, 0.5, size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    return slide

def score_slide():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, GRAY_BG)
    add_rect(slide, 0, 0, 13.33, 1.2, TEAL)
    add_rect(slide, 0, 1.2, 0.06, 6.3, ORANGE)
    add_text(slide, "Deliverable Coverage — Bobot Penilaian", 0.3, 0.2, 12, 0.85, size=26, bold=True, color=WHITE)
    rows = [
        ("Bagian 1 — Backend API",     "40%", "~90%", "✅ 10 endpoints, semua business rules"),
        ("Bagian 2 — Frontend Web App","40%", "~92%", "✅ 3 halaman + login + responsive"),
        ("Bagian 3 — Testing / README","15%", "~93%", "✅ 5 skenario test + 14 unit tests"),
        ("Bagian 4 — Refleksi & Tools","5%",  "100%", "✅ Tools, rejected output, arsitektur"),
    ]
    headers = ["Bagian", "Bobot", "Estimasi", "Status"]
    col_x = [0.3, 4.5, 5.8, 7.2]
    col_w = [4.1, 1.1, 1.2, 5.8]
    # header row
    for j, h in enumerate(headers):
        add_rect(slide, col_x[j], 1.4, col_w[j]-0.05, 0.45, TEAL_DARK)
        add_text(slide, h, col_x[j]+0.05, 1.45, col_w[j]-0.1, 0.38, size=12, bold=True, color=WHITE)
    for i, row in enumerate(rows):
        bg = WHITE if i % 2 == 0 else RGBColor(0xE8,0xF4,0xF4)
        y = 1.9 + i * 0.7
        for j, cell in enumerate(row):
            add_rect(slide, col_x[j], y, col_w[j]-0.05, 0.62, bg)
            color = TEAL_DARK if j == 0 else (GREEN if j == 1 else GRAY_TEXT)
            add_text(slide, cell, col_x[j]+0.05, y+0.1, col_w[j]-0.1, 0.48, size=12, bold=(j==0), color=color)
    add_rect(slide, 0.3, 4.7, 12.7, 0.55, TEAL_DARK)
    add_text(slide, "🏆  Total Estimasi Score: ~92 / 100", 0.4, 4.77, 12.5, 0.42, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return slide

# ════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════

# 1 — Cover
title_slide(
    "Mobile Sales Force (MSF)",
    "Case Study — Software Engineer | PT Mersifarma"
)

# 2 — Agenda
content_slide("Agenda Presentasi", [
    "1 — Latar Belakang & Objektif",
    "2 — Arsitektur Solusi",
    "3 — Business Process Flow",
    "4 — Backend API — Endpoint & Rules",
    "5 — Frontend Web App — Halaman & Fitur",
    "6 — Business Rules & Testing",
    "7 — Tools & Refleksi Pengerjaan",
    "8 — Demo & Live Run",
    "9 — Deliverable Coverage & Score",
    "10 — Tantangan & Future Development",
])

# 3 — Latar Belakang
section_slide("01 — Latar Belakang & Objektif")

content_slide("Latar Belakang & Objektif", [
    "##Konteks Bisnis",
    "PT Mersifarma membutuhkan sistem digital untuk Medical Representative (MR) mengelola kunjungan dokter",
    "Proses manual: Call List → approval → Call Plan → realisasi kunjungan",
    "##Objektif",
    "Bangun Backend API + Frontend Web App untuk 3 fitur utama",
    "Call List Bulanan → Approval Workflow → Call Plan → Call Actual",
    "##Target User",
    "MR (Medical Representative) — user utama, akses dari mobile",
    "DM (District Manager), RSM (Regional Sales Manager), MM (Marketing Manager) — approver",
])

# 4 — Arsitektur
section_slide("02 — Arsitektur Solusi")

two_col_slide(
    "Arsitektur — Monorepo Full-Stack",
    "🖥  Backend (Node.js + Express + PostgreSQL)",
    [
        "REST API di port 3000",
        "Auth: Bearer token (hardcoded, no JWT lib)",
        "Middleware: authenticate + requireRole",
        "Routes: /api/mcl, /api/products,",
        "  /api/call-lists, /api/call-plans,",
        "  /api/call-actuals",
        "DB: PostgreSQL 18, 8 tabel relasional",
        "Parameterized queries (SQL-injection safe)",
        "Schema: schema.sql | Seed: seed.sql + seed.js",
    ],
    "🌐  Frontend (React + Vite)",
    [
        "Dev server di port 5173",
        "Vite proxy /api → localhost:3000",
        "State: AuthContext + ToastContext",
        "Axios client dengan request/response interceptors",
        "Custom hooks: useCallLists, useCallPlans,",
        "  useCallActuals, useMCL, useProducts",
        "Pages: Login, CallList, CallPlan, CallActual",
        "CSS: mobile-first, design system tokens",
        "Tests: 14 unit tests (Vitest + RTL)",
    ]
)

# 5 — Business Flow
section_slide("03 — Business Process Flow")

flow_slide("End-to-End Business Flow", [
    ("MR\nLogin",      "Pilih akun\n(5 hardcoded\nusers)"),
    ("Buat\nCall List","Pilih bulan\n+ multi-select\ndokter dari MCL"),
    ("Submit\nCL",     "Status:\ndraft →\nsubmitted"),
    ("DM\nApprove",    "Atasan langsung\napprove / reject\n(single-level)"),
    ("Buat\nCall Plan","Dari CL approved:\ndokter + tanggal\n+ waktu"),
    ("Call\nActual",   "Terencana /\nUnplan /\nNon Target\n+ foto + TTD\n+ produk"),
])

# 6 — Approval Hierarchy
content_slide("Approval Hierarchy — Single Level", [
    "##Aturan Bisnis",
    "Setiap Call List hanya bisa di-approve oleh atasan langsung satu level di atas pengaju",
    "MR di-approve oleh DM | DM di-approve oleh RSM | RSM di-approve oleh MM",
    "RSM / MM tidak bisa approve Call List milik MR langsung → HTTP 403 Forbidden",
    "##Status Flow",
    "draft  →  submitted  →  approved   (happy path)",
    "draft  →  submitted  →  rejected   → kembali ke draft (MR bisa revisi & resubmit)",
    "##Implementasi",
    "Fungsi isDirectSupervisor() memetakan: { dm: 'mr', rsm: 'dm', mm: 'rsm' }",
    "Dipanggil pada setiap PATCH /api/call-lists/:id/approve sebelum proses",
])

# 7 — Backend
section_slide("04 — Backend API")

two_col_slide(
    "Backend API — Endpoint Coverage",
    "📋  Call List",
    [
        "POST   /api/call-lists — buat/upsert CL",
        "GET    /api/call-lists — list by role",
        "GET    /api/call-lists/:id — detail + dokter",
        "PATCH  /api/call-lists/:id/submit",
        "PATCH  /api/call-lists/:id/approve",
        "",
        "📅  Call Plan",
        "POST   /api/call-plans",
        "GET    /api/call-plans",
    ],
    "✅  Call Actual & Master Data",
    [
        "POST   /api/call-actuals",
        "GET    /api/call-actuals",
        "",
        "GET    /api/mcl   (master customers)",
        "GET    /api/products",
        "",
        "Error Codes:",
        "422 — validasi gagal",
        "409 — duplikat (visit/plan yang sama)",
        "403 — role tidak diizinkan",
        "401 — token tidak valid",
    ]
)

content_slide("Call Actual — 3 Mode Kunjungan", [
    "##Mode Terencana (plan_id diisi)",
    "doctor_id harus cocok dengan Call Plan → HTTP 422 jika tidak match",
    "Visit type otomatis: 'plan'",
    "##Mode Unplan (plan_id null, dokter ada di CL approved bulan ini)",
    "Cek call_list_doctors JOIN call_lists WHERE status='approved' & bulan sama",
    "Visit type otomatis: 'unplan'",
    "##Mode Non Target (plan_id null, dokter tidak ada di CL bulan ini)",
    "Cek master_customers — dokter harus terdaftar di MCL",
    "Visit type otomatis: 'non_target'",
    "##Rules Umum",
    "photo_url & signature_url wajib ada → 422",
    "Minimal 1 produk detailing → 422",
    "Maksimal 1 kunjungan per dokter per hari (UNIQUE constraint) → 409",
])

# 8 — Frontend
section_slide("05 — Frontend Web App")

two_col_slide(
    "Frontend — 3 Halaman Utama",
    "📋  Call List Page",
    [
        "Form: input bulan (type=month)",
        "Multi-select dokter dari MCL (checkbox)",
        "Tombol Simpan → POST /api/call-lists",
        "Tombol Submit → PATCH .../submit",
        "Tabel: ID, Bulan, Status badge, Dokter, Aksi",
        "Supervisor: tombol Setujui & Tolak",
        "Reject dialog dengan input alasan wajib",
    ],
    "📅  Call Plan & ✅ Call Actual",
    [
        "Call Plan: dropdown CL approved",
        "→ fetch dokter dari /api/call-lists/:id",
        "→ input tanggal + waktu kunjungan",
        "Call Actual: 3-mode tabs (Terencana/",
        "Unplan/Non Target)",
        "→ Terencana: dropdown Call Plan,",
        "   doctor auto-fill read-only",
        "→ Unplan/Non Target: dropdown MCL",
        "→ URL foto + URL tanda tangan",
        "→ Checkbox multi-select produk",
    ]
)

content_slide("Frontend — Quality & UX", [
    "##Integrasi API",
    "Axios client dengan request interceptor (inject Bearer token dari localStorage)",
    "Response interceptor: 401 → clear session + redirect /login otomatis",
    "##Loading & Error Handling",
    "Spinner pada semua operasi async (fetch + submit)",
    "Toast notification (auto-dismiss 3 detik) untuk sukses & error",
    "Error spesifik: 409 Duplikat, 422 Validasi, 403 Role — pesan berbeda tiap kasus",
    "##Responsive & Accessible",
    "Mobile-first CSS — layout berubah dari sidebar ke single column di layar kecil",
    "Semantic HTML: <main>, <section>, <table>, <form>, <button>",
    "ARIA labels, focus ring, min touch target 44×44px untuk MR di lapangan",
    "##Authentication",
    "5-card login (Andi MR, Sari MR, Doni DM, Citra RSM, Budi MM)",
    "Token disimpan di localStorage, role-aware NavBar & page redirects",
])

# 9 — Testing
section_slide("06 — Business Rules & Testing")

content_slide("5 Skenario Test Manual (README)", [
    "##Skenario 1 — Happy Path (End-to-End)",
    "MR buat CL → submit → DM approve → buat CP → catat CA → HTTP 201",
    "##Skenario 2 — Call Plan dari CL Belum Approved",
    "POST /api/call-plans dengan call_list_id berstatus 'draft' → HTTP 422",
    "Message: 'Call list must be approved to create a call plan'",
    "##Skenario 3 — Double Visit (Dokter Sama, Hari Sama)",
    "POST /api/call-actuals ke dokter yang sama dua kali → HTTP 409",
    "Message: 'A call actual for this doctor on this date already exists'",
    "##Skenario 4 — Call Actual Tanpa Foto",
    "POST /api/call-actuals tanpa photo_url → HTTP 422",
    "##Skenario 5 — Approve oleh Role Salah",
    "RSM atau MM approve Call List milik MR langsung → HTTP 403 Forbidden",
])

content_slide("Unit Tests — 14 Tests Passing", [
    "##axiosClient.test.js (5 tests)",
    "Request interceptor: token ada → Authorization header disertakan",
    "Request interceptor: tidak ada token → header tidak ditambah",
    "Response interceptor: 401 → localStorage clear + redirect /login",
    "Response interceptor: non-401 → error diteruskan, localStorage aman",
    "##LoginPage.test.jsx (4 tests)",
    "Render 5 kartu user | klik Andi → login + navigate /call-lists",
    "klik Budi (MM) → login dengan role mm | already logged in → redirect",
    "##PrivateRoute.test.jsx (2 tests)",
    "Unauthenticated → redirect /login | Authenticated → render children",
    "##CallActualForm.test.jsx (3 tests)",
    "Default mode Terencana: CP dropdown + doctor read-only",
    "Switch ke Unplan: CP hilang, MCL doctor dropdown muncul",
    "Switch ke Non Target: MCL doctor dropdown muncul",
])

# 10 — Tools
section_slide("07 — Tools & Refleksi")

two_col_slide(
    "Tools yang Digunakan",
    "🛠  Development Tools",
    [
        "Antigravity (AI coding assistant)",
        "→ scaffolding, SQL schema, unit test",
        "PostgreSQL 18 + psql CLI",
        "→ database relasional",
        "Node.js 22 + Express 4",
        "→ REST API backend",
        "React 18 + Vite 5",
        "→ frontend dev server + build",
        "Vitest + React Testing Library",
        "→ unit testing komponen React",
        "curl / Postman",
        "→ pengujian manual endpoint",
    ],
    "✅ Diterima vs ❌ Ditolak",
    [
        "✅ DITERIMA:",
        "isDirectSupervisor() map — logika tepat,",
        "  diverifikasi 6 kombinasi role",
        "resolveVisitType() — 3-branch logic,",
        "  parameterized query aman",
        "",
        "❌ DITOLAK:",
        "Rejection → status 'rejected' (bukan reset",
        "  ke draft) — bertentangan req 5.5",
        "String interpolation di SQL → SQL injection",
        "  risk, diganti parameterized query $1..$n",
        "FormData upload foto → backend JSON-only,",
        "  reverted ke URL text input",
    ]
)

# 11 — Score
score_slide()

# 12 — Challenges
section_slide("08 — Tantangan & Future Development")

two_col_slide(
    "Tantangan yang Dipertimbangkan",
    "🔴 Tantangan Teknis",
    [
        "Double visit via berbagai jalur",
        "→ UNIQUE(user_id, doctor_id, visit_date)",
        "   di call_actuals — satu constraint",
        "   mencegah semua jalur (plan/unplan/NT)",
        "",
        "CL direject setelah CP dibuat",
        "→ Data CP tetap ada (tidak auto-delete)",
        "→ Asumsi: CP tetap valid, perlu UI notif",
        "",
        "Visit_date di luar bulan CL",
        "→ Validasi bulan di POST /api/call-plans",
        "   menolak tanggal di luar periode CL",
    ],
    "🚀 Future Development",
    [
        "✅ Unplan Visit — implemented",
        "✅ Non Target Visit — implemented",
        "✅ Auth sederhana (login by role) — done",
        "✅ Unit test min 3 fungsi kritis — done",
        "",
        "⬜ Offline mode (localStorage sync)",
        "⬜ Pagination & filter pada list",
        "⬜ Deploy ke cloud (Railway/Render)",
        "⬜ GPS tracking check-in",
        "⬜ Reminder/escalation jika DM delay",
        "⬜ PATCH /call-actuals/:id/complete",
    ]
)

# 13 — Demo
section_slide("09 — Live Demo")

content_slide("Cara Menjalankan Aplikasi", [
    "##Prasyarat",
    "Node.js 18+  |  PostgreSQL (running)  |  psql CLI",
    "##Backend Setup",
    "1.  cp backend/.env.example backend/.env  → isi DATABASE_URL",
    "2.  psql -f backend/db/schema.sql",
    "3.  psql -f backend/db/seed.sql   (atau: node backend/db/seed.js)",
    "4.  node backend/server.js   → http://localhost:3000",
    "##Frontend Setup",
    "5.  cd frontend && npm install && npm run dev   → http://localhost:5173",
    "##Login",
    "Buka localhost:5173 → pilih salah satu dari 5 kartu user",
    "Token: token-mr1 (Andi MR) | token-dm (Doni DM) | token-mm (Budi MM)",
    "##Test",
    "cd backend && npm test   |   cd frontend && npm test",
])

# 14 — Closing
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, TEAL_DARK)
add_rect(slide, 0, 0, 13.33, 0.08, ORANGE)
add_rect(slide, 0, 7.42, 13.33, 0.08, ORANGE)
add_text(slide, "Terima Kasih", 0.5, 2.0, 12, 1.2, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Repository GitHub tersedia — siap di-review & di-demo", 0.5, 3.5, 12, 0.6, size=18, color=RGBColor(0xCC,0xEE,0xFF), align=PP_ALIGN.CENTER)
add_text(slide, "Q & A", 0.5, 4.4, 12, 0.7, size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

OUT = "/home/yusuf/Documents/Mobile Sales Force/MSF_Presentation.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
